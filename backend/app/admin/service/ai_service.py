#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import asyncio
import json
import re
import requests
from abc import ABC, abstractmethod
from typing import Optional, Dict, Any

from backend.common.log import log
from backend.app.admin.service.config_service import config_service
# 有可选标签时：优先从列表中选择，也可生成新标签
_CLASSIFY_SYSTEM_PROMPT = (
    "你是一个文本分类助手。根据用户提供的文本内容，优先从给定的标签列表中选出所有匹配的标签；"
    "若现有标签无法准确描述文本内容，可额外补充 1~3 个新标签，新标签应简洁（2~8 个字）且与现有标签风格一致。"
    "只返回一个 JSON 数组，数组元素为标签名称字符串，不要包含任何解释或多余文字。"
    "如果没有任何匹配或值得新增的标签，返回空数组 []。"
)

# 无可选标签时：由 AI 自由生成
_GENERATE_SYSTEM_PROMPT = (
    "你是一个文档标签助手。根据用户提供的文本内容，为其生成 3~5 个简洁的中文标签。"
    "标签应准确描述文档的主题、领域或内容类型，每个标签 2~8 个字，避免过于宽泛。"
    "只返回一个 JSON 数组，数组元素为标签名称字符串，不要包含任何解释或多余文字。"
)


# ---------------------------------------------------------------------------
# LLM 适配器
# ---------------------------------------------------------------------------

class LLMAdapter(ABC):
    """LLM适配器基类"""

    @abstractmethod
    def call(self, user_prompt: str, system_prompt: Optional[str] = None,
             config: Optional[Dict] = None) -> str:
        pass


class OpenAIAdapter(LLMAdapter):
    """OpenAI及兼容格式的适配器"""

    def __init__(self, api_key: str, base_url: str, model: str):
        self.api_key = api_key
        self.base_url = base_url
        self.model = model

    def call(self, user_prompt: str, system_prompt: Optional[str] = None,
             config: Optional[Dict] = None) -> str:
        max_tokens = 1000
        temperature = 0.2

        if config:
            llm_config = config.get("llm", {})
            max_tokens = llm_config.get("max_tokens", max_tokens)
            temperature = llm_config.get("temperature", temperature)

        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": user_prompt})

        payload = {
            'model': self.model,
            # 'max_tokens': max_tokens,
            'temperature': temperature,
            'messages': messages
        }

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}"
        }

        response = requests.post(self.base_url, headers=headers, json=payload)

        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
        else:
            raise Exception(f"API request failed: {response.text}")


class ClaudeAdapter(LLMAdapter):
    """Claude (Anthropic) 适配器"""

    def __init__(self, api_key: str, model: str = "claude-3-5-sonnet-20241022"):
        self.api_key = api_key
        self.model = model
        self.base_url = "https://api.anthropic.com/v1/messages"

    def call(self, user_prompt: str, system_prompt: Optional[str] = None,
             config: Optional[Dict] = None) -> str:
        max_tokens = 1000
        temperature = 0.2

        if config:
            llm_config = config.get("llm", {})
            max_tokens = llm_config.get("max_tokens", max_tokens)
            temperature = llm_config.get("temperature", temperature)

        payload = {
            'model': self.model,
            'max_tokens': max_tokens,
            'temperature': temperature,
            'messages': [{"role": "user", "content": user_prompt}]
        }

        if system_prompt:
            payload['system'] = system_prompt

        headers = {
            "Content-Type": "application/json",
            "x-api-key": self.api_key,
            "anthropic-version": "2023-06-01"
        }

        response = requests.post(self.base_url, headers=headers, json=payload)

        if response.status_code == 200:
            return response.json()['content'][0]['text']
        else:
            raise Exception(f"API request failed: {response.text}")


class GeminiAdapter(LLMAdapter):
    """Google Gemini 适配器"""

    def __init__(self, api_key: str, model: str = "gemini-pro"):
        self.api_key = api_key
        self.model = model
        self.base_url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    def call(self, user_prompt: str, system_prompt: Optional[str] = None,
             config: Optional[Dict] = None) -> str:
        temperature = 0.2

        if config:
            llm_config = config.get("llm", {})
            temperature = llm_config.get("temperature", temperature)

        # Gemini 将 system prompt 合并到 user prompt
        full_prompt = user_prompt
        if system_prompt:
            full_prompt = f"{system_prompt}\n\n{user_prompt}"

        payload = {
            "contents": [{
                "parts": [{"text": full_prompt}]
            }],
            "generationConfig": {
                "temperature": temperature
            }
        }

        response = requests.post(
            f"{self.base_url}?key={self.api_key}",
            headers={"Content-Type": "application/json"},
            json=payload
        )

        if response.status_code == 200:
            return response.json()['candidates'][0]['content']['parts'][0]['text']
        else:
            raise Exception(f"API request failed: {response.text}")


class LLMFactory:
    """LLM 工厂类"""

    @staticmethod
    def create_adapter(provider: str, **kwargs) -> LLMAdapter:
        """
        根据提供商创建适配器

        Args:
            provider: 提供商名称 (openai, claude, gemini, etc.)
            **kwargs: 各提供商所需的参数
        """
        if provider.lower() in ['openai', 'zhipu', 'deepseek', 'moonshot', 'yi', 'ollama']:
            return OpenAIAdapter(
                api_key=kwargs['api_key'],
                base_url=kwargs['base_url'],
                model=kwargs['model']
            )
        elif provider.lower() == 'claude':
            return ClaudeAdapter(
                api_key=kwargs['api_key'],
                model=kwargs.get('model', 'claude-3-5-sonnet-20241022')
            )
        elif provider.lower() == 'gemini':
            return GeminiAdapter(
                api_key=kwargs['api_key'],
                model=kwargs.get('model', 'gemini-pro')
            )
        else:
            raise ValueError(f"Unsupported provider: {provider}")


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

def extract_json_from_text(text):
    """
    Extract JSON array from text that might contain additional content.

    Args:
        text: Text that may contain JSON

    Returns:
        The parsed JSON if found, None otherwise
    """
    # First, check if the text is wrapped in code blocks with triple backticks
    code_block_pattern = r'```(?:json)?\s*([\s\S]*?)```'
    code_match = re.search(code_block_pattern, text)
    if code_match:
        text = code_match.group(1).strip()
        print("Found JSON in code block, extracting content...")

    try:
        # Try direct parsing in case the response is already clean JSON
        return json.loads(text)
    except json.JSONDecodeError:
        # Look for opening and closing brackets of a JSON array
        start_idx = text.find('[')
        if start_idx == -1:
            print("No JSON array start found in text")
            return None

        # Simple bracket counting to find matching closing bracket
        bracket_count = 0
        complete_json = False
        for i in range(start_idx, len(text)):
            if text[i] == '[':
                bracket_count += 1
            elif text[i] == ']':
                bracket_count -= 1
                if bracket_count == 0:
                    # Found the matching closing bracket
                    json_str = text[start_idx:i+1]
                    complete_json = True
                    break

        # Handle complete JSON array
        if complete_json:
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                print("Found JSON-like structure but couldn't parse it.")
                print("Trying to fix common formatting issues...")

                # Try to fix missing quotes around keys
                fixed_json = re.sub(r'(\s*)(\w+)(\s*):(\s*)', r'\1"\2"\3:\4', json_str)
                # Fix trailing commas
                fixed_json = re.sub(r',(\s*[\]}])', r'\1', fixed_json)

                try:
                    return json.loads(fixed_json)
                except:
                    print("Could not fix JSON format issues")
        else:
            # Handle incomplete JSON - try to complete it
            print("Found incomplete JSON array, attempting to complete it...")

            # Get all complete objects from the array
            objects = []
            obj_start = -1
            obj_end = -1
            brace_count = 0

            # First find all complete objects
            for i in range(start_idx + 1, len(text)):
                if text[i] == '{':
                    if brace_count == 0:
                        obj_start = i
                    brace_count += 1
                elif text[i] == '}':
                    brace_count -= 1
                    if brace_count == 0:
                        obj_end = i
                        objects.append(text[obj_start:obj_end+1])

            if objects:
                # Reconstruct a valid JSON array with complete objects
                reconstructed_json = "[\n" + ",\n".join(objects) + "\n]"
                try:
                    return json.loads(reconstructed_json)
                except json.JSONDecodeError:
                    print("Couldn't parse reconstructed JSON array.")
                    print("Trying to fix common formatting issues...")

                    # Try to fix missing quotes around keys
                    fixed_json = re.sub(r'(\s*)(\w+)(\s*):(\s*)', r'\1"\2"\3:\4', reconstructed_json)
                    # Fix trailing commas
                    fixed_json = re.sub(r',(\s*[\]}])', r'\1', fixed_json)

                    try:
                        return json.loads(fixed_json)
                    except:
                        print("Could not fix JSON format issues in reconstructed array")

        print("No complete JSON array could be extracted")
        return None


# ---------------------------------------------------------------------------
# 工具定义（OpenAI function calling 格式）
# ---------------------------------------------------------------------------

CHAT_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "search_docs",
            "description": "全文检索文档库，根据关键词搜索相关文件，返回匹配的文档列表（含 doc_id 和标题）",
            "parameters": {
                "type": "object",
                "properties": {
                    "keyword": {
                        "type": "string",
                        "description": "搜索关键词"
                    },
                    "size": {
                        "type": "integer",
                        "description": "返回结果数量，默认 20",
                        "default": 20
                    }
                },
                "required": ["keyword"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_doc_dir",
            "description": "创建一个新的目录，返回新目录的 dir_id",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "目录名称"
                    },
                    "parent_id": {
                        "type": "integer",
                        "description": "父目录 ID，不填则创建为顶级目录"
                    }
                },
                "required": ["name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "move_docs_to_dir",
            "description": "将指定的文档移动到某个目录",
            "parameters": {
                "type": "object",
                "properties": {
                    "doc_ids": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "要移动的文档 ID 列表"
                    },
                    "dir_id": {
                        "type": "integer",
                        "description": "目标目录 ID"
                    }
                },
                "required": ["doc_ids", "dir_id"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_text_doc",
            "description": "创建一篇文本文档，内容保存到知识库并建立全文检索索引",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "文档标题"
                    },
                    "content": {
                        "type": "string",
                        "description": "文档正文内容"
                    },
                    "doc_dir_id": {
                        "type": "integer",
                        "description": "所属目录 ID，不填则放在根目录"
                    }
                },
                "required": ["title", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_spreadsheet",
            "description": "创建一张表格文档，支持多行多列结构化数据",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "表格标题"
                    },
                    "headers": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "列名列表，如 [\"姓名\", \"金额\", \"日期\"]"
                    },
                    "rows": {
                        "type": "array",
                        "items": {
                            "type": "array",
                            "description": "一行数据，顺序与 headers 对应"
                        },
                        "description": "数据行列表，每行为与 headers 顺序对应的值数组"
                    },
                    "doc_dir_id": {
                        "type": "integer",
                        "description": "所属目录 ID，不填则放在根目录"
                    }
                },
                "required": ["title", "headers", "rows"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "list_dirs",
            "description": "获取所有文档目录的树形结构，返回 id、名称、父目录、子目录列表。在做整理操作前应先调用此工具了解现有目录结构",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_doc_content",
            "description": "获取指定文档的详细信息，包括标题、类型、内容摘要、标签、所属目录",
            "parameters": {
                "type": "object",
                "properties": {
                    "doc_id": {
                        "type": "integer",
                        "description": "文档 ID"
                    }
                },
                "required": ["doc_id"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "find_unclassified_docs",
            "description": "查找尚未归入任何目录的文档，用于治理检查和批量归档",
            "parameters": {
                "type": "object",
                "properties": {
                    "size": {
                        "type": "integer",
                        "description": "返回数量上限，默认 20",
                        "default": 20
                    }
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "tag_doc",
            "description": "为文档设置标签（覆盖写，传入完整标签列表）",
            "parameters": {
                "type": "object",
                "properties": {
                    "doc_id": {
                        "type": "integer",
                        "description": "文档 ID"
                    },
                    "tags": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "标签名称列表，如 [\"财务\", \"2024\"]，会覆盖原有标签"
                    }
                },
                "required": ["doc_id", "tags"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "create_ppt",
            "description": "创建一份 PPT 演示文稿，每页包含标题和正文内容",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "PPT 文件标题"
                    },
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "幻灯片标题"
                                },
                                "content": {
                                    "type": "string",
                                    "description": "幻灯片正文内容"
                                }
                            },
                            "required": ["title", "content"]
                        },
                        "description": "幻灯片列表，每项包含 title 和 content"
                    },
                    "doc_dir_id": {
                        "type": "integer",
                        "description": "所属目录 ID，不填则放在根目录"
                    }
                },
                "required": ["title", "slides"]
            }
        }
    }
]


# ---------------------------------------------------------------------------
# 工具执行函数
# ---------------------------------------------------------------------------

async def _execute_search_docs(keyword: str, size: int = 20) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    result = await SysDocService.search(keyword=keyword, page=1, size=size)
    items = result.get("items", [])
    return {
        "total": result.get("total", 0),
        "docs": [
            {
                "doc_id": item["doc_id"],
                "doc_title": item.get("doc_title", ""),
                "doc_name": item.get("doc_name", "")
            }
            for item in items
        ]
    }


async def _execute_create_doc_dir(name: str, parent_id: int = None) -> dict:
    from backend.app.admin.service.doc_dir_service import doc_dir_service
    from backend.app.admin.schema.doc_dir import CreateDocDirParam
    from backend.app.admin.crud.crud_doc_dir import doc_dir_dao
    from backend.database.db_pg import async_db_session

    obj = CreateDocDirParam(name=name, parent_id=parent_id)
    await doc_dir_service.create(obj=obj)

    async with async_db_session() as db:
        created = await doc_dir_dao.get_by_name(db, name, parent_id)

    if created:
        return {"success": True, "dir_id": created.id, "name": created.name}
    return {"success": True, "dir_id": None, "name": name}


async def _execute_move_docs_to_dir(doc_ids: list, dir_id: int) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import UpdateSysDocParam

    success_ids = []
    failed_ids = []
    for doc_id in doc_ids:
        try:
            await SysDocService.update(pk=doc_id, obj=UpdateSysDocParam(doc_dir_id=dir_id))
            success_ids.append(doc_id)
        except Exception as e:
            log.error(f"[move_docs_to_dir] 移动文档 {doc_id} 失败: {repr(e)}")
            failed_ids.append(doc_id)

    return {
        "success_count": len(success_ids),
        "failed_count": len(failed_ids),
        "success_ids": success_ids,
        "failed_ids": failed_ids
    }


async def _execute_create_text_doc(title: str, content: str, doc_dir_id: int = None) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import CreateSysDocParam

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        content=content,
        type="text",
        doc_dir_id=doc_dir_id,
    ))
    # 建立分块和全文检索向量
    await SysDocService.create_doc_tokens(id=doc.id)
    return {"success": True, "doc_id": doc.id, "title": doc.title}


async def _execute_create_spreadsheet(
    title: str,
    headers: list,
    rows: list,
    doc_dir_id: int = None,
) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import CreateSysDocParam
    from backend.app.admin.schema.doc_data import CreateSysDocDataParam
    from backend.app.admin.utils.tabular_processor import TabularProcessor

    # headers + rows → list[dict]
    data_json = [dict(zip(headers, row)) for row in rows]

    # 生成纯文本 content（供全文检索）
    content_lines = ["\t".join(str(v) for v in headers)]
    for row in rows:
        content_lines.append("\t".join(str(v) for v in row))
    content = "\n".join(content_lines)

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        content=content,
        type="spreadsheet",
        doc_dir_id=doc_dir_id,
    ))

    # 生成 Univer workbook JSON 并保存
    workbook_json = TabularProcessor.data_to_workbook(data_json, sheet_name=title)
    await SysDocService.base_update(pk=doc.id, obj={"workbook": workbook_json})

    # 保存行数据（供结构化查询）
    if data_json:
        obj_list = [CreateSysDocDataParam(doc_id=doc.id, row=r) for r in data_json]
        await SysDocService.create_doc_data(obj_list=obj_list)

    return {"success": True, "doc_id": doc.id, "title": doc.title, "row_count": len(rows)}


async def _execute_list_dirs() -> dict:
    from backend.app.admin.service.doc_dir_service import doc_dir_service

    tree = await doc_dir_service.get_doc_dir_tree()

    def _simplify(nodes: list) -> list:
        result = []
        for node in nodes:
            result.append({
                "dir_id": node.get("id"),
                "name": node.get("name"),
                "parent_id": node.get("parent_id"),
                "children": _simplify(node.get("children") or []),
            })
        return result

    return {"dirs": _simplify(tree)}


async def _execute_get_doc_content(doc_id: int) -> dict:
    from backend.app.admin.service.doc_service import sys_doc_service
    from backend.database.db_pg import async_db_session
    from sqlalchemy import select
    from backend.app.admin.model import SysDoc
    from sqlalchemy.orm import selectinload

    async with async_db_session() as db:
        result = await db.execute(
            select(SysDoc)
            .options(selectinload(SysDoc.tags))
            .where(SysDoc.id == doc_id)
        )
        doc = result.scalar_one_or_none()

    if not doc:
        return {"error": f"文档 {doc_id} 不存在"}

    content_preview = ""
    if doc.content:
        content_preview = doc.content[:500] + ("..." if len(doc.content) > 500 else "")

    return {
        "doc_id": doc.id,
        "title": doc.title,
        "type": doc.type,
        "doc_dir_id": doc.doc_dir_id,
        "tags": [t.name for t in (doc.tags or [])],
        "content_preview": content_preview,
        "created_time": str(doc.created_time) if doc.created_time else None,
    }


async def _execute_find_unclassified_docs(size: int = 20) -> dict:
    from backend.database.db_pg import async_db_session
    from sqlalchemy import select
    from backend.app.admin.model import SysDoc
    from sqlalchemy.orm import selectinload

    async with async_db_session() as db:
        result = await db.execute(
            select(SysDoc)
            .options(selectinload(SysDoc.tags))
            .where(SysDoc.doc_dir_id.is_(None))
            .order_by(SysDoc.created_time.desc())
            .limit(size)
        )
        docs = result.scalars().all()

    return {
        "total_shown": len(docs),
        "docs": [
            {
                "doc_id": d.id,
                "title": d.title,
                "type": d.type,
                "tags": [t.name for t in (d.tags or [])],
            }
            for d in docs
        ],
    }


async def _execute_tag_doc(doc_id: int, tags: list) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import UpdateSysDocParam

    await SysDocService.update(pk=doc_id, obj=UpdateSysDocParam(tags=tags))
    return {"success": True, "doc_id": doc_id, "tags": tags}


async def _execute_create_ppt(
    title: str,
    slides: list[dict],
    doc_dir_id: int = None,
) -> dict:
    import uuid
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.service.upload_service import UploadService
    from backend.app.admin.schema.doc import CreateSysDocParam
    from backend.core.conf import settings

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[1]  # 标题+内容布局

    for slide_data in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = slide_data.get("content", "")
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    ppt_bytes = buf.getvalue()

    unique_id = str(uuid.uuid4())
    file_suffix = ".pptx"
    new_filename = f"{unique_id}{file_suffix}"

    await UploadService._minio_put_object(
        settings.BUCKET_NAME,
        new_filename,
        ppt_bytes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    content_lines = []
    for s in slides:
        content_lines.append(s.get("title", ""))
        content_lines.append(s.get("content", ""))
    content = "\n".join(content_lines)

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        name=f"{title}.pptx",
        type="ppt",
        file=new_filename,
        uuid=unique_id,
        file_suffix=file_suffix,
        size=len(ppt_bytes),
        content=content,
        doc_dir_id=doc_dir_id,
        status=1,
    ))

    return {"success": True, "doc_id": doc.id, "title": doc.title, "slide_count": len(slides)}


async def _execute_tool(tool_name: str, tool_args: dict) -> str:
    try:
        if tool_name == "search_docs":
            result = await _execute_search_docs(**tool_args)
        elif tool_name == "create_doc_dir":
            result = await _execute_create_doc_dir(**tool_args)
        elif tool_name == "move_docs_to_dir":
            result = await _execute_move_docs_to_dir(**tool_args)
        elif tool_name == "create_text_doc":
            result = await _execute_create_text_doc(**tool_args)
        elif tool_name == "create_spreadsheet":
            result = await _execute_create_spreadsheet(**tool_args)
        elif tool_name == "list_dirs":
            result = await _execute_list_dirs()
        elif tool_name == "get_doc_content":
            result = await _execute_get_doc_content(**tool_args)
        elif tool_name == "find_unclassified_docs":
            result = await _execute_find_unclassified_docs(**tool_args)
        elif tool_name == "tag_doc":
            result = await _execute_tag_doc(**tool_args)
        elif tool_name == "create_ppt":
            result = await _execute_create_ppt(**tool_args)
        else:
            result = {"error": f"Unknown tool: {tool_name}"}
        return json.dumps(result, ensure_ascii=False)
    except Exception as e:
        log.error(f"[_execute_tool] 工具 {tool_name} 执行失败: {repr(e)}")
        return json.dumps({"error": repr(e)}, ensure_ascii=False)


# ---------------------------------------------------------------------------
# 统一 LLM 调用接口
# ---------------------------------------------------------------------------

async def call_llm(user_prompt: str, system_prompt: Optional[str] = None,
                   config: Optional[Dict] = None) -> str:
    """
    统一的 LLM 调用接口（异步版本）

    Args:
        user_prompt: 用户提示词
        system_prompt: 系统提示词
        config: 配置对象，包含 provider 和其他设置

    Returns:
        模型响应字符串
    """
    merged_settings = await config_service.get_merged_settings()
    provider = merged_settings.LLM_PROVIDER if merged_settings.LLM_PROVIDER else 'openai'

    adapter_kwargs = {
        'api_key': merged_settings.LLM_API_KEY,
        'base_url': merged_settings.LLM_API_URL,
        'model': merged_settings.LLM_MODEL
    }

    adapter = LLMFactory.create_adapter(provider, **adapter_kwargs)

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, adapter.call, user_prompt, system_prompt, config)


async def chat_with_tools(
    messages: list,
    system_prompt: Optional[str] = None,
    config: Optional[Dict] = None,
    max_tool_rounds: int = 10
) -> str:
    """
    带工具调用的对话接口（仅支持 OpenAI 兼容格式）。

    工具调用循环：LLM 返回 tool_calls → 执行工具 → 将结果追加到消息 → 再次调用 LLM，
    直到 LLM 返回普通文本或达到 max_tool_rounds 上限。

    Args:
        messages: 对话消息列表，格式同 OpenAI messages（role/content）
        system_prompt: 系统提示词
        config: 可选配置，支持 {"llm": {"temperature": 0.2}}
        max_tool_rounds: 最大工具调用轮次，防止无限循环

    Returns:
        LLM 最终的文本回复
    """
    merged_settings = await config_service.get_merged_settings()

    temperature = 0.2
    if config:
        temperature = config.get("llm", {}).get("temperature", temperature)

    all_messages = []
    if system_prompt:
        all_messages.append({"role": "system", "content": system_prompt})
    all_messages.extend(messages)

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {merged_settings.LLM_API_KEY}"
    }

    loop = asyncio.get_running_loop()

    for round_idx in range(max_tool_rounds):
        payload = {
            "model": merged_settings.LLM_MODEL,
            "temperature": temperature,
            "messages": all_messages,
            "tools": CHAT_TOOLS,
            "tool_choice": "auto"
        }

        response = await loop.run_in_executor(
            None,
            lambda p=payload: requests.post(merged_settings.LLM_API_URL, headers=headers, json=p)
        )

        if response.status_code != 200:
            raise Exception(f"API request failed: {response.text}")

        resp_data = response.json()
        choice = resp_data["choices"][0]
        assistant_message = choice["message"]
        finish_reason = choice.get("finish_reason", "stop")

        all_messages.append(assistant_message)

        if finish_reason == "tool_calls" and assistant_message.get("tool_calls"):
            for tool_call in assistant_message["tool_calls"]:
                tool_name = tool_call["function"]["name"]
                tool_args = json.loads(tool_call["function"]["arguments"])

                log.info(f"[chat_with_tools] round={round_idx+1} 调用工具: {tool_name}, 参数: {tool_args}")
                tool_result = await _execute_tool(tool_name, tool_args)
                log.info(f"[chat_with_tools] 工具结果: {tool_result}")

                all_messages.append({
                    "role": "tool",
                    "tool_call_id": tool_call["id"],
                    "content": tool_result
                })
        else:
            return assistant_message.get("content") or ""

    raise Exception(f"chat_with_tools 超出最大工具调用轮次 ({max_tool_rounds})")


# ---------------------------------------------------------------------------
# 业务方法
# ---------------------------------------------------------------------------

async def classify_text_tags(text: str, max_length: int = 2000) -> list[str]:
    """
    使用 LLM 对文本进行标签分类。

    - 数据库中有标签时：从现有标签中选择匹配的
    - 数据库中无标签时：由 AI 自由生成 3~5 个标签

    :param text: 输入文本
    :param max_length: 文本截取的最大长度
    :return: 标签名称列表
    """
    if not text or len(text.strip()) < 10:
        return []

    text_sample = text[:max_length] if len(text) > max_length else text

    # 读取数据库中已有的所有标签
    from backend.app.admin.crud.crud_tag import tag_dao
    from backend.database.db_pg import async_db_session

    existing_names: list[str] = []
    try:
        async with async_db_session() as db:
            all_tags = await tag_dao.get_all(db)
            existing_names = [t.name for t in all_tags if t.name]
    except Exception as e:
        log.warning(f"[classify_text_tags] 获取标签列表失败: {repr(e)}")

    if existing_names:
        tag_list_str = "、".join(existing_names)
        user_prompt = (
            f"现有标签：{tag_list_str}\n\n"
            f"文本内容：\n{text_sample}\n\n"
            f"请从现有标签中选出所有与文本相关的标签，若现有标签不足以描述文本内容，可额外补充新标签。以 JSON 数组格式返回所有标签名称。"
        )
        system_prompt = _CLASSIFY_SYSTEM_PROMPT
    else:
        user_prompt = (
            f"文本内容：\n{text_sample}\n\n"
            f"请为以上文本生成 3~5 个简洁的标签，以 JSON 数组格式返回。"
        )
        system_prompt = _GENERATE_SYSTEM_PROMPT

    try:
        response = await call_llm(user_prompt, system_prompt=system_prompt)

        start = response.find("[")
        end = response.rfind("]")
        if start == -1 or end == -1:
            log.warning(f"[classify_text_tags] LLM 返回内容无法解析为列表: {response!r}")
            return []

        raw_tags: list = json.loads(response[start : end + 1])

        # 接受所有非空字符串：既包含现有标签，也包含 AI 新生成的标签
        valid_tags = [t.strip() for t in raw_tags if isinstance(t, str) and t.strip()]

        if valid_tags:
            log.info(f"[classify_text_tags] AI 标签结果: {valid_tags}")

        return valid_tags

    except Exception as e:
        log.error(f"[classify_text_tags] AI 标签分类失败: {repr(e)}")
        return []


class AiService:
    call_llm = staticmethod(call_llm)
    chat_with_tools = staticmethod(chat_with_tools)
    classify_text_tags = staticmethod(classify_text_tags)
    extract_json_from_text = staticmethod(extract_json_from_text)


ai_service = AiService()

"""
文件操作工具执行函数

从 ai_service.py 迁出，接口改为统一 kwargs 调用。
"""
import uuid
import io

from backend.common.log import log


async def execute_search_docs(keyword: str, size: int = 20) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    result = await SysDocService.search(keyword=keyword, page=1, size=size)
    items = result.get("items", [])
    return {
        "total": result.get("total", 0),
        "docs": [
            {
                "doc_id": item["doc_id"],
                "doc_title": item.get("doc_title", ""),
                "doc_name": item.get("doc_name", ""),
            }
            for item in items
        ],
    }


async def execute_create_doc_dir(name: str, parent_id: int = None) -> dict:
    from backend.app.admin.service.doc_dir_service import doc_dir_service
    from backend.app.admin.schema.doc_dir import CreateDocDirParam
    from backend.app.admin.crud.crud_doc_dir import doc_dir_dao
    from backend.database.db_pg import async_db_session

    obj = CreateDocDirParam(name=name, parent_id=parent_id)
    await doc_dir_service.create(obj=obj)

    async with async_db_session() as db:
        created = await doc_dir_dao.get_by_name(db, name, parent_id)

    if created:
        return {"success": True, "dir_id": created.id, "name": created.name}
    return {"success": True, "dir_id": None, "name": name}


async def execute_move_docs_to_dir(doc_ids: list, dir_id: int) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import UpdateSysDocParam

    success_ids = []
    failed_ids = []
    for doc_id in doc_ids:
        try:
            await SysDocService.update(pk=doc_id, obj=UpdateSysDocParam(doc_dir_id=dir_id))
            success_ids.append(doc_id)
        except Exception as e:
            log.error(f"[move_docs_to_dir] 移动文档 {doc_id} 失败: {repr(e)}")
            failed_ids.append(doc_id)

    return {
        "success_count": len(success_ids),
        "failed_count": len(failed_ids),
        "success_ids": success_ids,
        "failed_ids": failed_ids,
    }


async def execute_create_text_doc(title: str, content: str, doc_dir_id: int = None) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import CreateSysDocParam

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        content=content,
        type="text",
        doc_dir_id=doc_dir_id,
    ))
    await SysDocService.create_doc_tokens(id=doc.id)
    return {"success": True, "doc_id": doc.id, "title": doc.title}


async def execute_create_spreadsheet(
    title: str,
    headers: list,
    rows: list,
    doc_dir_id: int = None,
) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import CreateSysDocParam
    from backend.app.admin.schema.doc_data import CreateSysDocDataParam
    from backend.app.admin.utils.tabular_processor import TabularProcessor

    data_json = [dict(zip(headers, row)) for row in rows]

    content_lines = ["\t".join(str(v) for v in headers)]
    for row in rows:
        content_lines.append("\t".join(str(v) for v in row))
    content = "\n".join(content_lines)

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        content=content,
        type="spreadsheet",
        doc_dir_id=doc_dir_id,
    ))

    workbook_json = TabularProcessor.data_to_workbook(data_json, sheet_name=title)
    await SysDocService.base_update(pk=doc.id, obj={"workbook": workbook_json})

    if data_json:
        obj_list = [CreateSysDocDataParam(doc_id=doc.id, row=r) for r in data_json]
        await SysDocService.create_doc_data(obj_list=obj_list)

    return {"success": True, "doc_id": doc.id, "title": doc.title, "row_count": len(rows)}


async def execute_list_dirs() -> dict:
    from backend.app.admin.service.doc_dir_service import doc_dir_service

    tree = await doc_dir_service.get_doc_dir_tree()

    def _simplify(nodes: list) -> list:
        result = []
        for node in nodes:
            result.append({
                "dir_id": node.get("id"),
                "name": node.get("name"),
                "parent_id": node.get("parent_id"),
                "children": _simplify(node.get("children") or []),
            })
        return result

    return {"dirs": _simplify(tree)}


async def execute_get_doc_content(doc_id: int) -> dict:
    from backend.database.db_pg import async_db_session
    from sqlalchemy import select
    from sqlalchemy.orm import selectinload
    from backend.app.admin.model import SysDoc

    async with async_db_session() as db:
        result = await db.execute(
            select(SysDoc)
            .options(selectinload(SysDoc.tags))
            .where(SysDoc.id == doc_id)
        )
        doc = result.scalar_one_or_none()

    if not doc:
        return {"error": f"文档 {doc_id} 不存在"}

    content_preview = ""
    if doc.content:
        content_preview = doc.content[:500] + ("..." if len(doc.content) > 500 else "")

    return {
        "doc_id": doc.id,
        "title": doc.title,
        "type": doc.type,
        "doc_dir_id": doc.doc_dir_id,
        "tags": [t.name for t in (doc.tags or [])],
        "content_preview": content_preview,
        "created_time": str(doc.created_time) if doc.created_time else None,
    }


async def execute_find_unclassified_docs(size: int = 20) -> dict:
    from backend.database.db_pg import async_db_session
    from sqlalchemy import select
    from sqlalchemy.orm import selectinload
    from backend.app.admin.model import SysDoc

    async with async_db_session() as db:
        result = await db.execute(
            select(SysDoc)
            .options(selectinload(SysDoc.tags))
            .where(SysDoc.doc_dir_id.is_(None))
            .order_by(SysDoc.created_time.desc())
            .limit(size)
        )
        docs = result.scalars().all()

    return {
        "total_shown": len(docs),
        "docs": [
            {
                "doc_id": d.id,
                "title": d.title,
                "type": d.type,
                "tags": [t.name for t in (d.tags or [])],
            }
            for d in docs
        ],
    }


async def execute_tag_doc(doc_id: int, tags: list) -> dict:
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.schema.doc import UpdateSysDocParam

    await SysDocService.update(pk=doc_id, obj=UpdateSysDocParam(tags=tags))
    return {"success": True, "doc_id": doc_id, "tags": tags}


async def execute_create_ppt(
    title: str,
    slides: list[dict],
    doc_dir_id: int = None,
) -> dict:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from backend.app.admin.service.doc_service import SysDocService
    from backend.app.admin.service.upload_service import UploadService
    from backend.app.admin.schema.doc import CreateSysDocParam
    from backend.core.conf import settings

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[1]

    for slide_data in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = slide_data.get("content", "")
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    ppt_bytes = buf.getvalue()

    unique_id = str(uuid.uuid4())
    file_suffix = ".pptx"
    new_filename = f"{unique_id}{file_suffix}"

    await UploadService._minio_put_object(
        settings.BUCKET_NAME,
        new_filename,
        ppt_bytes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    content_lines = []
    for s in slides:
        content_lines.append(s.get("title", ""))
        content_lines.append(s.get("content", ""))
    content = "\n".join(content_lines)

    doc = await SysDocService.create(obj=CreateSysDocParam(
        title=title,
        name=f"{title}.pptx",
        type="ppt",
        file=new_filename,
        uuid=unique_id,
        file_suffix=file_suffix,
        size=len(ppt_bytes),
        content=content,
        doc_dir_id=doc_dir_id,
        status=1,
    ))

    return {"success": True, "doc_id": doc.id, "title": doc.title, "slide_count": len(slides)}
